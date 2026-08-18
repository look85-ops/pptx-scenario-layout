"""
AI-пайплайн для генерации сценариев курсов.
Вход: техдокументация (docx) + инструкция методологии
Выход: структурированный сценарий дня в Markdown

Использование:
  python scripts/course_scenario_pipeline.py --module М4 --day 1 --sources "путь\к\файлу.docx"
"""

import os, sys, json, re, argparse
sys.stdout.reconfigure(encoding='utf-8')

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def extract_docx(path):
    """Извлечь текст из docx."""
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    tables_text = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                tables_text.append(" | ".join(cells))
    full = "\n".join(paragraphs)
    if tables_text:
        full += "\n\n=== ТАБЛИЦЫ ===\n" + "\n".join(tables_text)
    return full

def extract_pptx_text(path):
    """Извлечь текст из pptx (заглушка — возвращает имена слайдов)."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"--- Слайд {i} ---\n" + "\n".join(slide_texts))
        return "\n".join(texts)
    except ImportError:
        return f"[python-pptx не установлен, файл {path} пропущен]"
    except Exception as e:
        return f"[Ошибка чтения {path}: {e}]"

def build_prompt(module, day, sources_text, instruction_text, focus_notes=""):
    """Собрать промпт для LLM."""
    prompt = f"""Ты — педагогический дизайнер программ корпоративного обучения для руководителей РЖД.

Контекст программы: «Академия пути» — программа развития для резерва начальников службы пути (роль П — дорожный уровень). Аудитория: действующие руководители путевого комплекса с опытом 15+ лет.

Методологическая структура (обязательна):
Каждый блок строится по схеме «Вызов — Осмысление — Присвоение»:
1. Вызов (интерактивный вход): вопрос/ситуация из опыта аудитории, фиксация мнений
2. Осмысление (мини-лекция): подача теории тезисно, визуализация, связка с реальностью
3. Присвоение (практикум): групповая работа, защита решений, рефлексия

Тон общения: партнёрский, без менторства. Профессиональная лексика. Работа со скепсисом аудитории.

---
ЗАДАЧА: Разработать сценарий Модуль {module}, День {day}.

Исходные материалы (техдокументация):
{sources_text[:15000]}

{focus_notes}

---
ФОРМАТ ВЫВОДА (строго Markdown):

# Модуль {module} · День {day} · [Название дня]

Аудитория: резерв начальников службы пути (П)
Цель дня: [одним предложением]
Проверяемые результаты:
- [что участник сможет сделать после дня]

## Тайминг дня (общий)
| Время | Блок | Формат | Длительность |
|---|---|---|---|

## Блок 1: [Название] ([время])
### Цель блока
### Вызов (вопросы на вход, интерактив)
### Осмысление (тезисы, слайды)
### Присвоение (практикум, групповая работа)
### Быстрая проверка (1-2 вопроса)

## Блок 2: ...

## Итог дня (рефлексия, микроопрос)

ВАЖНО: Каждый слайд указывать в формате:
- ИД слайда (М{module}Д{day}_№XXX)
- Заголовок
- Основной текст
- Текст для преподавателя
- Интерактив/вопрос
"""
    return prompt

DEFAULT_INSTRUCTION = "INSTRUCTION_course.md"

def main():
    parser = argparse.ArgumentParser(description="AI-пайплайн сценариев курсов")
    parser.add_argument("--module", required=True, help="Номер модуля (М4)")
    parser.add_argument("--day", required=True, help="Номер дня (1-5)")
    parser.add_argument("--sources", nargs="+", help="Пути к исходным docx-файлам")
    parser.add_argument("--focus", default="", help="Доп. фокус для промпта")
    parser.add_argument("--instruction", default=DEFAULT_INSTRUCTION,
                        help=f"Путь к файлу инструкции методологии (по умолчанию: {DEFAULT_INSTRUCTION} в текущей папке)")
    parser.add_argument("--output", help="Куда сохранить результат (опционально)")
    parser.add_argument("--prompt-only", action="store_true", help="Только вывести промпт в stdout, без сохранения")
    parser.add_argument("--llm", nargs="?", const="deepseek/deepseek-chat",
                        help="Авто-генерация через LLM (указать 'provider/model'; env: LLM_API_KEY)")
    args = parser.parse_args()

    # 1. Read instruction
    instruction = ""
    if args.instruction and os.path.exists(args.instruction):
        with open(args.instruction, "r", encoding="utf-8") as f:
            instruction = f.read()
    elif args.instruction:
        print(f"Предупреждение: инструкция не найдена: {args.instruction}", file=sys.stderr)

    # 2. Extract sources
    sources_combined = []
    if args.sources:
        for path in args.sources:
            if not os.path.exists(path):
                sources_combined.append(f"=== {os.path.basename(path)} === ФАЙЛ НЕ НАЙДЕН\n")
                continue
            ext = os.path.splitext(path)[1].lower()
            if ext == ".docx":
                text = extract_docx(path)
            elif ext == ".pptx":
                text = extract_pptx_text(path)
            else:
                text = f"[Неподдерживаемый формат: {ext}]"
            sources_combined.append(f"=== {os.path.basename(path)} ({len(text)} chars) ===\n{text}\n")

    sources_text = "\n".join(sources_combined)

    # 3. Build prompt
    prompt = build_prompt(args.module, args.day, sources_text, instruction, args.focus)

    # 4. Output
    if args.prompt_only:
        print(prompt)
        print(f"--- Размер промпта: {len(prompt)} chars (~{len(prompt)//4} токенов) ---", file=sys.stderr)
        return

    output_path = args.output or os.path.join(BASE_DIR, "docs", "course", f"M{args.module}_D{args.day}_scenario.md")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    if args.llm:
        api_key = os.environ.get("LLM_API_KEY")
        api_url = os.environ.get("LLM_API_URL", "https://openai.bothub.chat/v1/chat/completions")
        if not api_key:
            print("Ошибка: задайте LLM_API_KEY в переменных окружения, ключ в код не кладём", file=sys.stderr)
            sys.exit(1)
        try:
            import requests
        except ImportError:
            print("Ошибка: для --llm нужен `requests` (pip install requests)", file=sys.stderr)
            sys.exit(1)
        resp = requests.post(
            api_url,
            headers={"Authorization": f"Bearer {api_key}"},
            json={
                "model": args.llm,
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.4,
                "max_tokens": 8000,
            },
            timeout=300,
        )
        resp.raise_for_status()
        text = resp.json()["choices"][0]["message"]["content"]
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"Сценарий сгенерирован и сохранён: {output_path}")
    else:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(prompt)
        print(f"Промпт сохранён: {output_path}")
        print(f"\nСкопируй этот промпт в Claude/ChatGPT для генерации сценария.")
        print(f"Или добавь --llm для авто-генерации (нужен LLM_API_KEY).")

if __name__ == "__main__":
    main()
